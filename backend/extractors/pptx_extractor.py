"""
PowerPoint sunumlarından metin çıkarır.
Slayt başlığı, içerik kutuları ve speaker notes ayrı ayrı toplanır;
notlar genelde dersin en yoğun anlatım kısmını içerir, bu yüzden atlanmaz.
"""

from pptx import Presentation

from .base import ExtractionError, ExtractionResult


def extract(path: str) -> ExtractionResult:
    try:
        prs = Presentation(path)
        parts: list[str] = []
        slide_count = 0

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_count += 1
            slide_parts = [f"--- Slayt {slide_idx} ---"]

            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide_parts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            slide_parts.append(" | ".join(cells))

            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_parts.append(f"[Sunum notu: {notes_text}]")

            if len(slide_parts) > 1:
                parts.extend(slide_parts)

        raw_text = "\n".join(parts)
        return ExtractionResult(raw_text=raw_text, metadata={"slide_count": slide_count})
    except Exception as e:
        raise ExtractionError(f"PPTX işlenemedi: {e}") from e
